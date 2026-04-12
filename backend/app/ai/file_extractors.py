"""Text extraction utilities for AI chat file attachments.

Extracts text content from various file types for inline inclusion
in LLM messages. Returns None for unsupported types or content
that exceeds the 512 KB limit.
"""

import io
import logging
from collections.abc import Callable

logger = logging.getLogger(__name__)

MAX_CONTENT_SIZE = 512 * 1024  # 512 KB


def extract_text(content: bytes, content_type: str) -> str | None:
    """Extract text from file content based on MIME type.

    Args:
        content: Raw file bytes.
        content_type: MIME type of the file.

    Returns:
        Extracted text string, or None for unsupported types
        or content exceeding the 512 KB limit.

    Raises:
        ValueError: If content exceeds the 512 KB size limit.
    """
    if len(content) > MAX_CONTENT_SIZE:
        raise ValueError(
            f"File content exceeds {MAX_CONTENT_SIZE // 1024} KB limit "
            f"({len(content)} bytes received)"
        )

    extractors: dict[str, Callable[[bytes], str]] = {
        "text/plain": _extract_passthrough,
        "text/csv": _extract_csv,
        "application/json": _extract_passthrough,
        "text/markdown": _extract_passthrough,
        "application/pdf": _extract_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
    }

    extractor = extractors.get(content_type)
    if extractor is None:
        return None

    try:
        return extractor(content)
    except Exception:
        logger.warning(
            "Failed to extract text from content_type=%s",
            content_type,
            exc_info=True,
        )
        return None


def _extract_passthrough(content: bytes) -> str:
    """Return raw UTF-8 decoded content for text types."""
    return content.decode("utf-8")


def _extract_csv(content: bytes) -> str:
    """Return CSV content as-is (text passthrough)."""
    return content.decode("utf-8")


def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    """Extract paragraph text from DOCX using python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_xlsx(content: bytes) -> str:
    """Extract cell data from XLSX as CSV-like text using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            rows.append(",".join(cells))
    wb.close()
    return "\n".join(rows)


def _extract_pptx(content: bytes) -> str:
    """Extract slide text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)
